"""Route each manifest-ingested file to the right extractor, persist raw text per-doc.

Outputs:
  work/extracted/<safe_filename>.json   # one file per source document
  work/skipped_docs.json                 # files with <200 chars extracted (flagged for OCR)
"""
from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import sys
import traceback
import unicodedata
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Callable, List, Optional

from dotenv import load_dotenv
from tqdm import tqdm

from pipeline.manifest import ManifestEntry, ingest_worklist


load_dotenv()

REPO_ROOT = Path(__file__).resolve().parent.parent
WORK_DIR = REPO_ROOT / "work"
EXTRACT_DIR = WORK_DIR / "extracted"
SKIPPED_LOG = WORK_DIR / "skipped_docs.json"

SOURCE_ROOT = Path(
    os.path.expanduser(
        os.getenv("SOURCE_CONTENT_PATH", "~/Desktop/Claude/synexis-bot/source_content")
    )
)

# Manifest paths for the baseline corpus (Patents/, Published Studies/, Manuals and
# Guides/, Regulatory Opinions/, Sales and Marketing Materials/, Public Domain
# Materials ex OSHA EPA etc/, and the root-level SDS/technical bulletins) strip
# the "Customer Facing Training Documents/" prefix that exists on disk. Probe both.
PATH_FALLBACK_PREFIXES: tuple[str, ...] = ("", "Customer Facing Training Documents/")

MIN_CHARS_FOR_INGEST = 200


# Many of the corpus filenames came from Windows/Word authoring and contain
# smart punctuation (curly quotes, en/em dashes) or non-breaking spaces, while
# the manifest was hand-typed with plain ASCII. Normalize both sides before
# comparing so paths resolve despite these invisible-to-the-eye mismatches.
def _normalize_for_match(s: str) -> str:
    s = unicodedata.normalize("NFC", s)
    s = (
        s.replace("\u2018", "'")
        .replace("\u2019", "'")
        .replace("\u201c", '"')
        .replace("\u201d", '"')
        .replace("\u2013", "-")
        .replace("\u2014", "-")
        .replace("\xa0", " ")
    )
    return re.sub(r"\s+", " ", s).strip()


def _resolve_source(relative_path: str) -> Optional[Path]:
    for prefix in PATH_FALLBACK_PREFIXES:
        p = SOURCE_ROOT / (prefix + relative_path)
        if p.exists():
            return p

    # Fuzzy fallback: scan the expected parent dir for a filename whose
    # punctuation-normalized form matches the manifest name.
    rel = Path(relative_path)
    target_name_norm = _normalize_for_match(rel.name)
    for prefix in PATH_FALLBACK_PREFIXES:
        parent = SOURCE_ROOT / (prefix + str(rel.parent))
        if not parent.is_dir():
            continue
        for candidate in parent.iterdir():
            if _normalize_for_match(candidate.name) == target_name_norm:
                return candidate
    return None


@dataclass
class PageText:
    number: int           # 1-based page or slide number
    text: str


@dataclass
class ExtractedDoc:
    file_path: str                 # relative path from SOURCE_ROOT (matches manifest)
    doc_id: str                    # deterministic hash of file_path
    source_category: str
    intake_mode: str
    description: str
    extension: str
    extractor_used: str
    total_chars: int
    pages: List[PageText] = field(default_factory=list)


def _doc_id(relative_path: str) -> str:
    return hashlib.sha1(relative_path.encode("utf-8")).hexdigest()[:16]


def _safe_name(relative_path: str) -> str:
    return re.sub(r"[^A-Za-z0-9_.-]+", "_", relative_path)


# ---------- per-extension extractors ----------

def _extract_pdf(abs_path: Path) -> tuple[str, List[PageText]]:
    """Try pdfplumber first; fall back to pymupdf if it yields too little text."""
    import pdfplumber  # type: ignore

    pages: List[PageText] = []
    try:
        with pdfplumber.open(str(abs_path)) as pdf:
            for i, p in enumerate(pdf.pages, start=1):
                txt = p.extract_text() or ""
                pages.append(PageText(number=i, text=txt))
    except Exception:
        pages = []
    total = sum(len(p.text) for p in pages)
    if total >= MIN_CHARS_FOR_INGEST:
        return "pdfplumber", pages

    # fallback
    import fitz  # pymupdf
    pages2: List[PageText] = []
    try:
        with fitz.open(str(abs_path)) as doc:
            for i, page in enumerate(doc, start=1):
                pages2.append(PageText(number=i, text=page.get_text("text") or ""))
    except Exception:
        pages2 = []
    total2 = sum(len(p.text) for p in pages2)
    if total2 > total:
        return "pymupdf", pages2
    return "pdfplumber", pages


def _extract_docx(abs_path: Path) -> tuple[str, List[PageText]]:
    """python-docx has no page concept — treat the whole doc as page 1."""
    from docx import Document

    doc = Document(str(abs_path))
    parts: List[str] = []
    for para in doc.paragraphs:
        if para.text:
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text for c in row.cells if c.text]
            if cells:
                parts.append(" | ".join(cells))
    text = "\n".join(parts).strip()
    return "python-docx", [PageText(number=1, text=text)]


def _extract_pptx(abs_path: Path) -> tuple[str, List[PageText]]:
    from pptx import Presentation

    prs = Presentation(str(abs_path))
    pages: List[PageText] = []
    for i, slide in enumerate(prs.slides, start=1):
        chunks: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    runs = "".join(run.text for run in para.runs)
                    if runs.strip():
                        chunks.append(runs)
            if shape.shape_type == 19 or getattr(shape, "has_table", False):
                # table
                try:
                    tbl = shape.table
                    for row in tbl.rows:
                        cells = [c.text for c in row.cells if c.text]
                        if cells:
                            chunks.append(" | ".join(cells))
                except Exception:
                    pass
        # speaker notes
        if slide.has_notes_slide:
            notes_txt = slide.notes_slide.notes_text_frame.text if slide.notes_slide.notes_text_frame else ""
            if notes_txt.strip():
                chunks.append(f"[NOTES] {notes_txt}")
        pages.append(PageText(number=i, text="\n".join(chunks).strip()))
    return "python-pptx", pages


EXTRACTORS: dict[str, Callable[[Path], tuple[str, List[PageText]]]] = {
    "pdf": _extract_pdf,
    "docx": _extract_docx,
    "pptx": _extract_pptx,
}


# ---------- driver ----------

def extract_one(entry: ManifestEntry) -> tuple[Optional[ExtractedDoc], Optional[dict]]:
    abs_path = _resolve_source(entry.relative_path)
    if abs_path is None:
        return None, {
            "file_path": entry.relative_path,
            "reason": "file_not_found",
            "tried_prefixes": list(PATH_FALLBACK_PREFIXES),
        }

    fn = EXTRACTORS.get(entry.extension)
    if fn is None:
        return None, {
            "file_path": entry.relative_path,
            "reason": f"no_extractor_for_extension:{entry.extension}",
        }

    try:
        extractor_used, pages = fn(abs_path)
    except Exception as exc:
        return None, {
            "file_path": entry.relative_path,
            "reason": "extraction_error",
            "error": f"{type(exc).__name__}: {exc}",
            "traceback": traceback.format_exc(limit=3),
        }

    total_chars = sum(len(p.text) for p in pages)
    if total_chars < MIN_CHARS_FOR_INGEST:
        return None, {
            "file_path": entry.relative_path,
            "reason": "near_empty",
            "extractor_used": extractor_used,
            "total_chars": total_chars,
            "needs_ocr": True,
        }

    doc = ExtractedDoc(
        file_path=entry.relative_path,
        doc_id=_doc_id(entry.relative_path),
        source_category=entry.source_category,
        intake_mode=entry.intake_mode,
        description=entry.description,
        extension=entry.extension,
        extractor_used=extractor_used,
        total_chars=total_chars,
        pages=pages,
    )
    return doc, None


def _write_doc(doc: ExtractedDoc) -> Path:
    EXTRACT_DIR.mkdir(parents=True, exist_ok=True)
    out_path = EXTRACT_DIR / f"{doc.doc_id}__{_safe_name(doc.file_path)}.json"
    payload = asdict(doc)
    out_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return out_path


def run(limit: Optional[int] = None, only: Optional[List[str]] = None) -> dict:
    WORK_DIR.mkdir(parents=True, exist_ok=True)
    entries = ingest_worklist()
    if only:
        only_set = set(only)
        entries = [e for e in entries if e.relative_path in only_set]
    if limit is not None:
        entries = entries[:limit]

    extracted_count = 0
    skipped: List[dict] = []

    for entry in tqdm(entries, desc="extract"):
        doc, skip = extract_one(entry)
        if doc is not None:
            _write_doc(doc)
            extracted_count += 1
        if skip is not None:
            skipped.append(skip)

    SKIPPED_LOG.write_text(json.dumps(skipped, ensure_ascii=False, indent=2), encoding="utf-8")
    summary = {
        "entries_considered": len(entries),
        "extracted": extracted_count,
        "skipped": len(skipped),
        "extracted_dir": str(EXTRACT_DIR),
        "skipped_log": str(SKIPPED_LOG),
    }
    print(json.dumps(summary, indent=2))
    return summary


def main(argv: List[str] | None = None) -> int:
    ap = argparse.ArgumentParser(description="Extract text from manifest-ingested sources.")
    ap.add_argument("--limit", type=int, default=None, help="Process only the first N manifest entries.")
    ap.add_argument(
        "--files",
        type=str,
        default=None,
        help="Comma-separated relative paths (must match manifest) to process instead of the full worklist.",
    )
    args = ap.parse_args(argv)
    only = [s.strip() for s in args.files.split(",")] if args.files else None
    run(limit=args.limit, only=only)
    return 0


if __name__ == "__main__":
    sys.exit(main())
